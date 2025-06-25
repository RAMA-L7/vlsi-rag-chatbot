import os
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.utils import embedding_functions

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

# Load embedding model
print("🔄 Loading embedding model...")
model = SentenceTransformer("all-MiniLM-L6-v2")

# Collect text from various file types
def load_text_from_file(filepath):
    if filepath.endswith(".txt"):
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    elif filepath.endswith(".pdf"):
        doc = fitz.open(filepath)
        return "\n".join([page.get_text() for page in doc])
    elif filepath.endswith(".docx"):
        from docx import Document
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    elif filepath.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    else:
        return ""


# Read and chunk all documents
print("📂 Reading files from /data...")
docs = []
for filename in os.listdir("data"):
    full_path = os.path.join("data", filename)
    text = load_text_from_file(full_path)
    if text.strip():
        chunks = text.split("\n\n")  # basic paragraph split
        docs.extend(chunks)

print(f"📄 Loaded {len(docs)} chunks.")

# Setup Chroma DB
print("📦 Setting up Chroma DB...")
chroma_client = chromadb.PersistentClient(path="./chromadb")
# If collection exists, delete it first
existing = chroma_client.list_collections()
if any(c.name == "vlsi_docs" for c in existing):
    chroma_client.delete_collection("vlsi_docs")

collection = chroma_client.create_collection("vlsi_docs")


# Embed and store
print("✨ Embedding and storing...")
for i, chunk in enumerate(docs):
    if chunk.strip():
        embedding = model.encode(chunk).tolist()
        collection.add(
            documents=[chunk],
            embeddings=[embedding],
            ids=[f"doc-{i}"]
        )

print("✅ Done! All chunks embedded and stored.")
